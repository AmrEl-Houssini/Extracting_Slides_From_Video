import cv2
import os
import zipfile
import io
from pptx import Presentation
import streamlit as st

# Function to extract unique slides
def extract_unique_slides(video_path, threshold=30):
    video_capture = cv2.VideoCapture(video_path)
    success, frame = video_capture.read()
    if not success:
        st.error("Failed to open video.")
        return []

    previous_frame = None
    slide_images = []
    slide_count = 0

    st.info("Extracting slides screenshots...")

    while success:
        current_frame_gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)

        if previous_frame is not None:
            frame_diff = cv2.absdiff(previous_frame, current_frame_gray)
            _, diff_thresh = cv2.threshold(frame_diff, threshold, 255, cv2.THRESH_BINARY)
            non_zero_count = cv2.countNonZero(diff_thresh)

            if non_zero_count > 1000:  # Adjust based on sensitivity needs
                slide_count += 1
                # Save slide to in-memory buffer
                is_success, buffer = cv2.imencode(".png", frame)
                if is_success:
                    slide_images.append(buffer)
        previous_frame = current_frame_gray
        success, frame = video_capture.read()

    video_capture.release()
    st.success("Slides extracted successfully.")
    return slide_images

# Function to create PowerPoint presentation
def create_presentation(slide_images):
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    st.info("Building the PowerPoint presentation...")

    for buffer in slide_images:
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        img_stream = io.BytesIO(buffer)
        slide.shapes.add_picture(img_stream, 0, 0, width=slide_width, height=slide_height)

    # Save presentation to in-memory buffer
    output_pptx = io.BytesIO()
    prs.save(output_pptx)
    output_pptx.seek(0)
    st.success("PowerPoint built successfully.")
    return output_pptx

# Function to create a ZIP file for slides
def create_zip(slide_images):
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "w") as zip_file:
        for idx, buffer in enumerate(slide_images, start=1):
            zip_file.writestr(f"slide_{idx}.png", buffer)
    zip_buffer.seek(0)
    return zip_buffer

# Streamlit app starts here
st.title("Educational Video Slide Extractor")

# Initialize session state for slide images and PowerPoint file
if "slide_images" not in st.session_state:
    st.session_state.slide_images = None
if "pptx_file" not in st.session_state:
    st.session_state.pptx_file = None

# Upload video file
uploaded_file = st.file_uploader("Upload Video File", type=["mp4", "avi", "mov"])
if uploaded_file is not None:
    st.success("Video uploaded successfully!")

    # Save uploaded video temporarily in memory
    temp_video_path = os.path.join("temp_video.mp4")
    with open(temp_video_path, "wb") as f:
        f.write(uploaded_file.read())

    # Extract slides
    if st.button("Start Slide Extraction"):
        st.session_state.slide_images = extract_unique_slides(temp_video_path)

    # Show download buttons after slides are extracted
    if st.session_state.slide_images:
        # Provide download button for slides screenshots as a ZIP
        zip_file = create_zip(st.session_state.slide_images)
        st.download_button(
            label="Download All Slides as ZIP",
            data=zip_file,
            file_name="slides.zip",
            mime="application/zip",
        )

        # Build PowerPoint if slides are extracted
        if st.button("Build PowerPoint"):
            st.session_state.pptx_file = create_presentation(st.session_state.slide_images)

        # Provide download link for the PowerPoint
        if st.session_state.pptx_file:
            st.download_button(
                label="Download PowerPoint Presentation",
                data=st.session_state.pptx_file,
                file_name="presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )